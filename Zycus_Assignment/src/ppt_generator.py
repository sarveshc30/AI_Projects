"""Builds the 5-7 slide monthly executive deck synthesizing both sample projects.

Per the brief: this stage reuses `key_risks` / `recommended_actions` from stage 5's LLM
output as slide-copy source material -- it does NOT make a fresh LLM call to "regenerate"
insights. Cross-project synthesis (slide 3, slide 4 ranking) is rule-based keyword grouping
over that reused material, run in this module.
"""
from __future__ import annotations

import re
import sys
from collections import defaultdict
from pathlib import Path
from typing import Optional

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from src.pipeline import PipelineResult, run_pipeline_for_file
from src.sentiment import SentimentAnalyzer

OUTPUT_PATH = Path(__file__).resolve().parent.parent / "outputs" / "deck" / "Monthly_Executive_Deck.pptx"
DATA_SAMPLES_DIR = Path(__file__).resolve().parent.parent / "data" / "samples"

RAG_COLORS = {
    "Green": RGBColor(0x1E, 0x8E, 0x3E),
    "Amber": RGBColor(0xE3, 0x8A, 0x00),
    "Red": RGBColor(0xC5, 0x22, 0x1F),
}

THEME_KEYWORDS = {
    "Milestone delays": ["milestone"],
    "Schedule slippage": ["schedule", "slippage", "behind", "delay"],
    "Stakeholder sentiment / communication": ["sentiment", "comment", "stakeholder"],
    "Blockers & dependencies": ["blocker", "predecessor", "on hold", "block"],
    "Data quality gaps": ["data gap", "missing", "absent", "unavailable", "no usable", "empty"],
}


def _short_name(project_display_name: str) -> str:
    # "Zycus - Titan S2P Implementation" -> "Titan S2P"; keeps slide copy compact.
    m = re.search(r"-\s*(.+?)\s+Implementation", project_display_name)
    return m.group(1) if m else project_display_name


def _classify_theme(text: str) -> Optional[str]:
    lowered = text.lower()
    for theme, keywords in THEME_KEYWORDS.items():
        if any(k in lowered for k in keywords):
            return theme
    return None


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_bullet_slide(prs: Presentation, title: str, bullets: list[str], note: Optional[str] = None) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
    if note:
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(9), Inches(0.6))
        tf = note_box.text_frame
        tf.text = note
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.italic = True


def _add_table_slide(
    prs: Presentation, title: str, headers: list[str], rows: list[list[str]], status_col: Optional[int] = None
) -> None:
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    n_rows, n_cols = len(rows) + 1, len(headers)
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(0.5 + 0.5 * n_rows)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(13)
            if status_col is not None and c == status_col and value in RAG_COLORS:
                cell.text_frame.paragraphs[0].font.color.rgb = RAG_COLORS[value]
                cell.text_frame.paragraphs[0].font.bold = True


def _portfolio_snapshot_rows(results: list[PipelineResult]) -> list[list[str]]:
    rows = []
    for r in results:
        driver = (
            r.rag_result.overrides_triggered[0]
            if r.rag_result.overrides_triggered
            else f"Weighted composite ({r.rag_result.composite_score}/100)"
        )
        rows.append([
            _short_name(r.project.project_display_name),
            r.rag_result.final_status,
            str(r.rag_result.composite_score),
            driver,
        ])
    return rows


def _build_cross_project_themes(results: list[PipelineResult]) -> dict[str, set[str]]:
    """theme -> set of short project names whose key_risks/data_gaps mention it."""
    theme_projects: dict[str, set[str]] = defaultdict(set)
    for r in results:
        name = _short_name(r.project.project_display_name)
        texts = list(r.rag_result.overrides_triggered)
        if r.llm_reasoning:
            texts += r.llm_reasoning.key_risks + r.llm_reasoning.data_gaps
        for t in texts:
            theme = _classify_theme(t)
            if theme:
                theme_projects[theme].add(name)
    return theme_projects


def _build_ranked_risks(results: list[PipelineResult]) -> list[tuple[str, list[str]]]:
    """Dedup similar risks by theme, tag with the project(s) that raised them, rank by breadth."""
    theme_risks: dict[str, list[str]] = defaultdict(list)
    theme_projects: dict[str, set[str]] = defaultdict(set)
    for r in results:
        name = _short_name(r.project.project_display_name)
        if not r.llm_reasoning:
            continue
        for risk in r.llm_reasoning.key_risks:
            theme = _classify_theme(risk) or "Other"
            theme_risks[theme].append(risk)
            theme_projects[theme].add(name)

    ranked = sorted(theme_projects.items(), key=lambda kv: (-len(kv[1]), kv[0]))
    output = []
    for theme, projects in ranked:
        tag = "Both projects" if len(projects) > 1 else next(iter(projects))
        example = theme_risks[theme][0]
        output.append((theme, [f"[{tag}] {example}"]))
    return output


def _build_recommendations(results: list[PipelineResult]) -> list[str]:
    seen_themes: set[str] = set()
    recs: list[str] = []
    for r in results:
        if not r.llm_reasoning:
            continue
        name = _short_name(r.project.project_display_name)
        for action in r.llm_reasoning.recommended_actions[:3]:
            theme = _classify_theme(action) or action[:30]
            if theme in seen_themes:
                continue
            seen_themes.add(theme)
            recs.append(f"[{name}] {action}")
    recs.append(
        "Portfolio-wide: stand up the weekly automated reporting cadence "
        "(config/schedule.yaml, docs/schedule.md) so RAG drift is caught between monthly reviews."
    )
    return recs


def generate_deck(
    file_paths: Optional[list[Path]] = None, output_path: Path = OUTPUT_PATH
) -> Path:
    file_paths = file_paths or sorted(DATA_SAMPLES_DIR.glob("*.xlsx"))
    analyzer = SentimentAnalyzer()
    results = [run_pipeline_for_file(f, sentiment_analyzer=analyzer, write_report=False) for f in file_paths]

    as_of_dates = [
        r.project.summary.get("todays_date") for r in results if r.project.summary.get("todays_date")
    ]
    period_str = max(as_of_dates).strftime("%B %Y") if as_of_dates else "Current period"

    prs = Presentation()

    # Slide 1: Title
    project_names = ", ".join(_short_name(r.project.project_display_name) for r in results)
    _add_title_slide(
        prs,
        "Project Health — Monthly Executive Review",
        f"Reporting period: {period_str}\nPortfolio: {project_names}",
    )

    # Slide 2: Portfolio RAG overview
    _add_table_slide(
        prs,
        "Portfolio RAG Overview",
        ["Project", "RAG Status", "Composite Score", "Primary Driver"],
        _portfolio_snapshot_rows(results),
        status_col=1,
    )

    # Slide 3: Cross-project trends & recurring themes
    theme_projects = _build_cross_project_themes(results)
    shared = [t for t, ps in theme_projects.items() if len(ps) > 1]
    unique = [(t, next(iter(ps))) for t, ps in theme_projects.items() if len(ps) == 1]
    theme_bullets = []
    if shared:
        theme_bullets.append("Recurring across both projects:")
        theme_bullets += [f"  - {t}" for t in shared]
    if unique:
        theme_bullets.append("Project-specific themes:")
        theme_bullets += [f"  - {t} ({p})" for t, p in unique]
    if not theme_bullets:
        theme_bullets = ["No recurring risk themes detected across the two projects this period."]
    _add_bullet_slide(
        prs,
        "Cross-Project Trends & Recurring Themes",
        theme_bullets,
        note=(
            "This is a first synthesis across two projects with one time-snapshot each -- "
            "these are cross-project patterns to date, not a fabricated multi-period trend line."
        ),
    )

    # Slide 4: Emerging risks, ranked, tagged
    ranked_risks = _build_ranked_risks(results)
    risk_bullets = [f"{theme}: {items[0]}" for theme, items in ranked_risks[:6]]
    if not risk_bullets:
        risk_bullets = ["No LLM-articulated risks available this period (see data gaps)."]
    _add_bullet_slide(prs, "Emerging Risks (Ranked)", risk_bullets)

    # Slide 5: Compact project snapshots
    snapshot_rows = []
    for r in results:
        pct = r.metrics.schedule.get("actual_pct_complete")
        pct_str = f"{pct * 100:.0f}%" if pct is not None else "N/A"
        top_risk = r.llm_reasoning.key_risks[0] if r.llm_reasoning and r.llm_reasoning.key_risks else "N/A"
        snapshot_rows.append([
            _short_name(r.project.project_display_name),
            r.rag_result.final_status,
            pct_str,
            top_risk,
        ])
    _add_table_slide(
        prs,
        "Compact Project Snapshots",
        ["Project", "RAG", "% Complete", "Top Risk"],
        snapshot_rows,
        status_col=1,
    )

    # Slide 6: Recommendations & next steps
    _add_bullet_slide(prs, "Recommendations & Next Steps", _build_recommendations(results))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


if __name__ == "__main__":
    path = generate_deck()
    print(f"Deck written to {path}")
